#pptx slide manipulation
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE

from shPaths import *
from shLib import getCsvData, logEntry
# from slDeck_single import singleDeck
# from slDeck_multi import multiDeck
# import slDeck_single as slSingle
# import slDeck_multi as slMulti

#from shLib import logEntry
import re 
import datetime
# import slDeck_single
# import slDeck_multi


class Table_Options:
    """ Sets either defaults or given options for a table"""
    def __init__(self):
        #set default options
        self.presentation = None

        self.custData = None
        self.archv_opt = None
        #!
        self.sanList = None # the list of SAN names and csvPath of retrieved csv files

        self.csvFile = None # current csv file name to extract data from
        #!
        self.csvColumns = None # current csv  columns (letters) to extract
        self.csvPivotCols = None
        self.csvPivotHeaders = None

        self.dbConnection = None # db Connection handler: conn = sql.connect(sqlite_file)
        self.dbTableName = None #name of the db table to create
        self.dbColNames = None #Column names to use on the dbtable [(names,...)]
        
        self.hBand = False #sets horizontal shade banding
        self.vBand = False #sets vertical shade banding
        self.font_size = Pt(18) #<<<--------CHANGE to just use an integer
        self.font_bold = False
        self.font_italic = False
        self.font_name = 'Calibri'
        self.font_color = RGBColor(12, 34, 56) # Black
        self.text_hAlign = PP_ALIGN.CENTER #LEFT, CENTER, RIGHT --Try literal conversion ast.eval() 
        self.text_vAlign = MSO_ANCHOR.MIDDLE #TOP, MIDDLE, BOTTOM -- to enter only a string options = 'right'
        self.first_row = True
        self.first_col = False
        self.last_row = False
        self.last_col = False
        self.left = Inches(1.0)
        self.top = Inches(2.25)
        #!
        self.title = ' '
        self.subtitle = ' '
        self.subtitle_fontSize = Pt(30)

        #TODO?
        #from pptx.enum.dml import MSO_THEME_COLOR
        #font.color.theme_color = MSO_THEME_COLOR.ACCENT_1

#---------------------------------------------------------------
def formatDbUsed(data):
    #convert the bytes value into a MB to shorten the width of the column
    newData = []
    #grab the dbUse value
    for item in data:
        row = tuple(item)
        #Example: (row[6] = 11.7% of 1045274B)
        #Example: (row[6] = 11% of 1045274B)
        # extract 1045274B and convert to 1.0MB
        #items = re.match('(\d+\.\d+%)\s+of\s+(\d+)B', row[5])
        #1) grab everything before the % sign
        #2) grab everything between 'of' and 'B'
        items = re.match('(.+)%\s+of(.+)B', row[5])
        #print items.group(1)
        #print items.group(2)

        # get the two items in parenthesis (the usage % and the db size in B)
        usage = round( float(items.group(1)), 1)
        # convert Bytes to MB
        mb = float(items.group(2))/1000000
        lst = list(row[:5]) # save the first 4 elements
        # add the last element(db usage) to the tuple
        #(tuples are immutable so first convert to a list 
        lst.append(str(usage) + '%' + '  of ' + str(round(mb, 1)) + ' MB')
        newData.append(tuple(lst))
    return newData
#---------------------------------------------------------------
##def formatFruStatus(data):
##    ''' format the fru status to show either none (for blank, ok or enabled)
##        or other'''
##    
##    newData=[]
##    #grab the swtich Status value
##    for row in data:
##
##        status = row[5].lower()
##        if (status == 'ok') or\
##           (status == 'enabled') or\
##           (status == ''):
##            status = 'None'
##
##        lst = list(row[:5]) # save the first 4 elements
##        # add the last element(db usage) to the tuple
##        #(tuples are immutable so first convert to a list 
##        lst.append(str(status))
##        newData.append(tuple(lst))
##    return newData

def fSize(numLines, maxFsize=20):    
    '''
    Returns the font size in points (an integer between 10-20)
        for the given number of lines, up to the maxFSize
    '''
        
    if numLines < 12:
        font_size = 20
        
    elif numLines < 13:
        font_size = 19
        
    elif numLines < 14:
        font_size = 18
        
    elif numLines < 15:
        font_size = 17
        
    elif numLines < 16:
        font_size = 16
        
    elif numLines < 19:
        font_size = 15
        
    elif numLines < 20:
        font_size = 14
        
    elif numLines < 22:
        font_size = 13
         
    elif numLines < 24:
        font_size = 12
        
    elif numLines < 26:
        font_size = 11
        
    else:
        font_size = 10
        
    if font_size > maxFsize:
        font_size = maxFsize
    return font_size

    

def format_table(table, options):
    def iter_cells(table):
        for row in table.rows:
            for cell in row.cells:
                yield cell

    #set given table options            
    table.horz_banding = options.hBand
    table.vert_banding = options.vBand
    table.last_col = options.last_col
    table.first_row = options.first_row
    table.last_row = options.last_row
    table.first_col = options.first_col
    table.last_col = options.last_col

    for cell in iter_cells(table):
        
        #the row height adjusts automatically to font size.
        #set margins as a table options?
        cell.margin_left = cell.margin_right = Pt(10)
        cell.margin_top = cell.margin_bottom = 0
        cell.vertical_anchor = options.text_vAlign
        
        for paragraph in cell.text_frame.paragraphs:
            paragraph.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
            run = cell.text_frame.paragraphs[0]

            # A run is each line of text in a paragraph
            # Just looking for line 0 because...
            # ...iterating over multiple runs in a cell disables horizontal alignment
            # (pptx bug?)
            run.font.size = options.font_size
            run.font.bold = options.font_bold
            run.font.italic = options.font_italic
            run.font.name = options.font_name
            run.alignment = options.text_hAlign
            #run.font.color.rgb = options.font_color

##------------------------------------------------------------------------------------VVVVVVVVV

def create_table(slide, data, options):
    title = options.title
    subtitle = options.subtitle
    #width = Inches(10) #total width of table
    width = Inches(1) #total width of table
    height = Inches(.1) #total height of table
    left = options.left
    top = options.top
    rows = len(data)
    # the first row --data[0]-- should always be the headers (see slides) 
    # (i.e., exact number of columns)
    #note that other rows might be group headers with only one column!
    cols = len(data[0]) 

    #create table
    shapes = slide.shapes
    shapes.title.text = title

    #Add subtitle
    #add_textbox(left, top, width, height)
    subtxt = shapes.add_textbox(Inches(1),Inches(1.5), Inches(10), Inches(0.5))
    subtxt.text = subtitle
    subtxt.text_frame.paragraphs[0].font.size = options.subtitle_fontSize

                                  
    table = shapes.add_table(rows, cols, left, top, width, height).table

    #input data to table's cells
    try:
        for row_idx,row in enumerate(data):
            count = 0 #count the numner of non-empty cells
            for col_idx, cell in enumerate(row):
                table.cell(row_idx, col_idx).text = str(cell)
                #count the number of columns in this row
                #if just one column --> group header row
                count += 1
            #Merge group header cells for rows with just one element
            #Check if all row cells (i.e., columns) are empty (optionally except the first)
            if count == 1: # header row            
                #mergeCellsHorizontally(table, row_idx, start_col_idx, end_col_idx)
                mergeCellsHorizontally(table, row_idx, 0, cols-1)

                #change row color
                table.cell(row_idx, 0).fill.solid()
                # set foreground (fill) color to a specific RGB color
                #table.cell(row_idx, 0).fill.fore_color.rgb = RGBColor(0xFB, 0x8F, 0x00) # Orange
                #table.cell(row_idx, 0).fill.fore_color.rgb = RGBColor(0x00, 0xFF, 0x00) #lima green
                table.cell(row_idx, 0).fill.fore_color.rgb = RGBColor(0x66, 0xFF, 0x66) #pastel green green
    except:
        print 'SLIDE ERROR!: Data columns exceed number of headers (slLib/create_table)'
        print 'Slide:', title, ' - Headers:', len(data[0]), data[0], 'Columns:', len(row), 'Row', row_idx+1
        print 'Quiting' 
        quit()

    # Manual way to set specific columns' widths
    #(if omited, all columns are equally sized)
    #table.columns[0].width = Inches(1.0)
    #table.columns[1].width = Inches(1.0)

    #Adjust column widths based on max word length from all cells in a column
    maxLen = []
    for row in xrange(len(data)):
        if len(data[row]) == 1:
            #divider header / single merged cell-- do not use for width calculation
            continue
        for col, value in enumerate(data[row]):
            #For the header column only, row == 0, 
            #break each string into space-separated words 
            #to find the longest word in the string.
            #Powerpoint will automatically break strings in cells 
            #into mutiple lines using the spaces (if any), 
            #if the cell is not wide enough.
            if row == 0:
                value = str(value)
                #For header row only, break the string into separate words
                #to measure their length individually. 
                words = value.split(' ')
            else:
                #For the rest of rows, make the column wide enough
                #to display the cell text on a single line
                words = [value]
            
            #find the longest word in the list
            for word in words:                    
                txtLen = len(str(word))
                try:
                    if txtLen > maxLen[col]:
                        maxLen[col] = txtLen
                except:
                    #maxLen[] is an empty list on first pass (header row) 
                    maxLen.append(txtLen)

    #format font size based on the number of lines
    #Optional max font size limt of 18pt: Tables with few rows look better.
    font_size = fSize(rows, 18)


    #---------------------------------
    #---------------------------------
    #FOR TESTING
    #font_size = 20
    #---------------------------------
    #---------------------------------


    
    options.font_size = Pt(font_size)


    #for each column, set the width BASED ON THE MAX LENGHT DATA STRING
    #the number of columns is the len() of any row. Using just first row, data[0].
    for col in xrange(len(data[0])):
        table.columns[col].width = colWidth(maxLen[col], font_size)
        
    # Manually set specific columns' widths
    #(if omited, all columns are equally sized)
    #table.columns[0].width = Inches(0)
    #table.columns[1].width = Inches(1.0)

       
    format_table(table, options)
    #reset font size in case it was changed
    options.subtitle_fontSize = Pt(30)
    return table
##------------------------------------------------------------------------------------^^^^^^^^^

def colWidth(max_word_length, font_size):
    '''
    Returns the column width based on the given length of the longest word
    and the font size in Pt.
    '''
        #Find the correction FACTOR
#!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!
        #NOTE:
        #AS THE FONT GETS SMALLER, THE CORRECTION BECOMES MORE NON-LINEAR
        #THE CORRECTION FACTOR NEEDS TO DECREASE 
        #AS THE NUMBER OF SYMBOLS INCREASES
        #otherwise, the smaller strings fit right in 
        #but the longer ones have too much empty space around.
    if font_size > 19:
        weight = 0.57
    elif font_size > 18:
        weight = 0.55
    elif font_size > 12:
        weight = 0.60
    elif font_size > 11:
        weight = 0.65
    elif font_size > 10:
        weight = 0.70
    else: weight = 1.00
#!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!
    
    
    return (Pt( font_size + font_size * max_word_length * weight ) )



def create_single_table(options):
    #Wrapper to create a slide with just one table from csv data

    #generate slide for table
    prs = options.presentation
    title_only_slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(title_only_slide_layout)
    
    #get data to populate table
    shData = getCsvData(options)

    #create table
    table = create_table(slide, shData, options)

    return table


def create_single_table_db(shData, options):
    #Wrapper to create a slide with just one table from db data

    #generate slide for table
    prs = options.presentation
    title_only_slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(title_only_slide_layout)
    
    #create table
    table = create_table(slide, shData, options)

    return table


# merge cells vertically
#https://groups.google.com/forum/#!topic/python-pptx/cVRP9sSpEjA
def mergeCellsVertically(table, start_row_idx, end_row_idx, col_idx):
    row_count = end_row_idx - start_row_idx + 1
    column_cells = [r.cells[col_idx] for r in table.rows][start_row_idx:]

    column_cells[0]._tc.set('rowSpan', str(row_count))
    for c in column_cells[1:]:
        c._tc.set('vMerge', '1')

# merge cells horizontally
#https://groups.google.com/forum/#!topic/python-pptx/cVRP9sSpEjA
def mergeCellsHorizontally(table, row_idx, start_col_idx, end_col_idx):
    col_count = end_col_idx - start_col_idx + 1
    row_cells = [c for c in table.rows[row_idx].cells][start_col_idx:end_col_idx]
    row_cells[0]._tc.set('gridSpan', str(col_count))
    for c in row_cells[1:]:
        c._tc.set('hMerge', '1')



def copy_slide(pres,pres1,index):
    source = pres.slides[index]
    #blank_slide_layout = _get_blank_slide_layout(pres)
    blank_slide_layout = pres.slide_layouts[blank_slide_layout_index]
    dest = pres1.slides.add_slide(blank_slide_layout)

    for shp in source.shapes:
        el = shp.element
        newel = copy.deepcopy(el)
        dest.shapes._spTree.insert_element_before(newel, 'p:extLst')

    # NOTE: six.iteritems = dictionary.iteritems() on Python 2 and dictionary.items() on Python 3.
    for key, value in six.iteritems(source.part.rels):
        # Make sure we don't copy a notesSlide relation as that won't exist
        if not "notesSlide" in value.reltype:
            dest.part.rels.add_relationship(value.reltype, value._target, value.rId)
    return dest


def move_slide(presentation, old_index, new_index):
        xml_slides = presentation.slides._sldIdLst  # pylint: disable=W0212
        slides = list(xml_slides)
        xml_slides.remove(slides[old_index])
        xml_slides.insert(new_index, slides[old_index])


def textSlide(presentation, title, subtitle, text, **opts):
    '''
    Creates a slide with a title, a subtitle and 
    multiple rows of text with options for text color 
    and textbox placement and size.
    **opts examples:
    
    #text Color (R, G, B)
    txt_color = (138,43,226)
    
    #text box (left, top, width, height)
    txt_box = (1.5, 2.75, 10, 3.5) 
    '''
    
    #default text color (138,43,226) # blueviolet
    txt_color = opts.get("txt_color", (138,43,226) )
    #textbox(left, top, width, height)
    #default txt_box = (1.5, 2.75, 10, 3.5)
    txt_box = opts.get("txt_box", (1.5, 2.75, 10, 3.5) )
    
    textbox =[]
    for par in txt_box:
        textbox.append(Inches(par))

    prs = presentation
    
    title_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(title_slide_layout)
    slide_title = slide.shapes.title
    slide_subtitle = slide.placeholders[1]
    slide_title.text = title
    slide_subtitle.text = subtitle
    
    #add_textbox(left, top, width, height)
    shape = slide.shapes
    
    #txtBox = shape.add_textbox(Inches(1.5),Inches(2.75), Inches(10), Inches(3.5))
    txtBox = shape.add_textbox(*textbox)
    
    font_size = fSize(len(text))
    #fSize = 17
    txtBox.text_frame.paragraphs[0].font.size = Pt(font_size)
    #txtBox.text_frame.paragraphs[0].font.color.rgb = RGBColor(138,43,226) # blueviolet
    txtBox.text_frame.paragraphs[0].font.color.rgb = RGBColor(*txt_color) # blueviolet
    #Add new-line characters (\n) between each text element        
    paragraph = '\n'.join(text)
    txtBox.text = paragraph


#--------------------------------------------------------------------
def addHeaders(headers, data):    
    '''
    Add column headers to print on the slide table
    'Headers' is a tuple with the column names added on 
    the very first record of the 'data' list
    '''
    if len(data) == 0:
        logEntry("No Slide Data!")
        return None
    if len(headers[0]) != len(data[0]):
        print 'SLIDE ERROR!: Mismatch between number of data columns and headers (slLib/addHeaders)'
        print 'Headers:', len(headers[0]), headers[0], ', Columns:', len(data[0])
        print 'Quiting' 
        quit()

    headers.extend(data)
    data = headers
    return data

#--------------------------------------------------------------------
def saveDeck(tbl_options):
    ''' Store a remote and local copy of the slide deck
    This is used by both, slDeck_single and slDeck_multi'''
    #close the connection to the database file
    #Note: the db may need to remain open if using the RAM file option
    #tbl_options.dbConnection.close()

    #!
    #save the slide deck to the customer's SH directory
    #using the current san health file name
    #customer, csvPath, shName, sanName, shYear = tbl_options.custData
    customer, csvPath, shName, sanName, shDate, shYear = tbl_options.custData
    folder = drive + startFolder + customer + shFolder
    #!
    
    #but if a slide deck with the same name already exists and it is open
    #add a timestamp to the name to make it unique    
    timestamp = datetime.datetime.now().strftime("%y-%m-%d-%H%M")
    datestamp = datetime.datetime.now().strftime("%y-%m-%d")

    prs = tbl_options.presentation

    if sanName == 'ALL':
        shName = customer + '_AGGREGATE_' + datestamp
    try:
        prs.save(folder + shName + '.pptx')
    except:
        #if slide deck with the same name is open...
        #... store a new one renamed with a timestamp 
        print folder + shName, 'Kept Open'
        prs.save(folder + shName + '-'+ timestamp + '.pptx')
    #save a LOCAL copy of the slide deck
    folder = archiveFolder
    try:
        prs.save(folder + shName + '.pptx')
    except:
        #if slide deck with the same name is open...
        #... store a new one renamed with a timestamp 
        print folder + shName, 'Kept Open'
        prs.save(folder + shName + '-'+ timestamp + '.pptx')
#-----------------------------------------------------------------

# #----------------------------------------------------------------------
# def createSlideDeck(tbl_options):
#     #This function creates a slide deck
#     #from all the SAN Reports downloaded
#  
#     #IF THERE IS NO DATA FOR A SLIDE, PRINT A NOTE ON THE SLIDE: NO DATA!
#     #--------------------------------------------------------
#     #SLIDE CREATION
#     #--------------------------------------------------------
#     # 1) loop over the options.sanList tuplets (sanName, csvPath),
#     # 2) store current san pointer into tbl_options.custData
#     # 3) call slide deck creator (single or multi)
# 
#     for san in tbl_options.sanList:
#         #retrieve next SAN data
#         customer, csvPath, shName, sanName, shYear = tbl_options.custData
#         shName, shFile, sanName, csvPath = san
#         #shName: John_Morrison_170726_1640_Maiden_Prod
#         #shFile: 7-27-2017_John_Morrison_170726_1640_Maiden_Prod.zip
#         
#         #and store it for the slide creator function
#         custData = (customer, csvPath, shName, sanName, shYear)
#         tbl_options.custData = custData
#         
#         #make and save slideDeck
#         slSingle.singleDeck(tbl_options)
#         #print 'SAN', sanName
#         logEntry('Slides Created', customer, shName)
#         
#     if len(tbl_options.sanList) > 1:
#         # create a deck with the agregated data from all the downloaded reports
#         #store multi SAN directive to use the customer name as the file name.
#         sanName = 'ALL'
#         custData = (customer, csvPath, shName, sanName, shYear)
#         tbl_options.custData = custData
# 
#         slMulti.multiDeck(tbl_options)
#         #print 'SAN', sanName
#         logEntry('Slides Created', customer, 'Agregate')
#    # END OF SLIDES
#    #note: the db connection is opened by loadDbTables
#     tbl_options.dbConnection.close()
# # #--------------------------------------------------------------------

